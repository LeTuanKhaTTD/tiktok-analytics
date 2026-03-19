"""Tao bao cao PowerPoint - Phan tich TikTok @travinhuniversity
Chi tap trung TikTok - khong co YouTube
18 slides: Title > Muc luc > Quy trinh > Thu thap > Xu ly & Gan nhan
         > KPI > Bar > Top10 > Line > Scatter
         > Sentiment > Chi tiet Sentiment > WordCloud > WC +-
         > Top Comments > Engagement x Sentiment > Tom tat > Thank you
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# Color scheme
C_PRIMARY = RGBColor(0x1B, 0x3A, 0x5C)
C_SECONDARY = RGBColor(0x2E, 0x86, 0xC1)
C_ACCENT = RGBColor(0xE7, 0x4C, 0x3C)
C_SUCCESS = RGBColor(0x27, 0xAE, 0x60)
C_WARNING = RGBColor(0xF3, 0x9C, 0x12)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT = RGBColor(0xEC, 0xF0, 0xF1)
C_DARK = RGBColor(0x2C, 0x3E, 0x50)
C_BLACK = RGBColor(0x33, 0x33, 0x33)
C_GRAY = RGBColor(0x99, 0x99, 0x99)
C_PURPLE = RGBColor(0x8E, 0x44, 0xAD)
C_TEAL = RGBColor(0x16, 0xA0, 0x85)

CHART_DIR = Path('reports/charts')


def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def txt(slide, l, t, w, h, text, sz=14, bold=False, color=C_BLACK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Arial'
    p.alignment = align
    return box


def bullets(slide, l, t, w, h, items, sz=13, color=C_DARK):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.name = 'Arial'
        p.space_after = Pt(6)
    return box


def table(slide, l, t, w, h, headers, rows, col_widths=None):
    shape = slide.shapes.add_table(len(rows)+1, len(headers),
                                    Inches(l), Inches(t), Inches(w), Inches(h))
    tbl = shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Inches(cw)
    for i, hdr in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = hdr
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = C_WHITE
            p.font.name = 'Arial'
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_PRIMARY
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r+1, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = C_BLACK
                p.font.name = 'Arial'
                p.alignment = PP_ALIGN.CENTER
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_LIGHT
    return shape


def img(slide, name, l, t, w, h):
    path = CHART_DIR / name
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(l), Inches(t), Inches(w), Inches(h))
        return True
    return False


def card(slide, x, y, w, h, value, label, color):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    txt(slide, x+0.1, y+0.15, w-0.2, 0.7, value, sz=26, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(slide, x+0.1, y+0.9, w-0.2, 0.4, label, sz=11, color=C_WHITE, align=PP_ALIGN.CENTER)


def step_box(slide, x, y, w, h, number, title, desc, color):
    """Draw a pipeline step box with number circle, title, and description."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    # Number
    txt(slide, x+0.1, y+0.1, 0.5, 0.5, number, sz=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    # Title
    txt(slide, x+0.6, y+0.1, w-0.8, 0.4, title, sz=14, bold=True, color=C_WHITE)
    # Desc
    txt(slide, x+0.6, y+0.5, w-0.8, h-0.6, desc, sz=10, color=RGBColor(0xDD,0xDD,0xDD))


def arrow_right(slide, x, y):
    """Draw a right arrow between steps."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                    Inches(x), Inches(y), Inches(0.5), Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_GRAY
    shape.line.fill.background()


def create():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # =================================================================
    # SLIDE 1: TITLE
    # =================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C_PRIMARY)
    txt(s, 1, 0.8, 11, 1.2, 'PHAN TICH DU LIEU TIKTOK', sz=42, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(s, 1, 2.2, 11, 0.8, '@travinhuniversity - Truong Dai hoc Tra Vinh', sz=24, color=C_LIGHT, align=PP_ALIGN.CENTER)

    # Divider line
    shape = s.shapes.add_shape(1, Inches(4), Inches(3.2), Inches(5.33), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_LIGHT
    shape.line.fill.background()

    txt(s, 1, 3.5, 11, 0.6, '100 Videos  |  1,222 Binh luan  |  10.8M Luot xem', sz=18, color=C_LIGHT, align=PP_ALIGN.CENTER)

    # Info box
    txt(s, 1, 4.6, 11, 0.5, 'Cong cu: Apify TikTok Scraper + PhoBERT Sentiment Analysis', sz=14, color=C_LIGHT, align=PP_ALIGN.CENTER)
    txt(s, 1, 5.1, 11, 0.5, 'Nguoi thuc hien: [Ten] | Bao cao tuan 03/03 - 10/03/2026', sz=14, color=C_LIGHT, align=PP_ALIGN.CENTER)
    txt(s, 1, 6.2, 11, 0.5, 'Truong Dai hoc Tra Vinh - Thuc tap', sz=13, color=C_GRAY, align=PP_ALIGN.CENTER)

    # =================================================================
    # SLIDE 2: MUC LUC
    # =================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C_WHITE)
    txt(s, 0.5, 0.3, 12, 0.7, 'MUC LUC BAO CAO', sz=28, bold=True, color=C_PRIMARY)

    sections = [
        ('PHAN 1', 'PHUONG PHAP THUC HIEN', 'Quy trinh pipeline, thu thap du lieu, xu ly & gan nhan cam xuc', C_SECONDARY, 'Slide 3-5'),
        ('PHAN 2', 'PHAN TICH TUONG TAC', 'Tong quan KPI, so sanh video, xu huong luot xem, tuong quan', C_SUCCESS, 'Slide 6-10'),
        ('PHAN 3', 'PHAN TICH CAM XUC', 'Sentiment analysis, confidence, word cloud, binh luan tieu bieu', C_PURPLE, 'Slide 11-15'),
        ('PHAN 4', 'KET HOP & KET LUAN', 'Engagement x Sentiment, tom tat, de xuat cai thien', C_ACCENT, 'Slide 16-18'),
    ]

    for i, (part, title, desc, clr, slides) in enumerate(sections):
        y = 1.3 + i * 1.4
        # Color bar
        shape = s.shapes.add_shape(1, Inches(0.5), Inches(y), Inches(0.15), Inches(1.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = clr
        shape.line.fill.background()
        # Part label
        txt(s, 0.9, y, 1.5, 0.4, part, sz=12, bold=True, color=clr)
        # Title
        txt(s, 0.9, y+0.3, 7, 0.4, title, sz=18, bold=True, color=C_DARK)
        # Desc
        txt(s, 0.9, y+0.7, 7, 0.4, desc, sz=12, color=C_GRAY)
        # Slide numbers
        txt(s, 9, y+0.3, 3, 0.4, slides, sz=14, bold=True, color=clr, align=PP_ALIGN.RIGHT)

    # Summary bar at bottom
    shape = s.shapes.add_shape(1, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_PRIMARY
    shape.line.fill.background()
    txt(s, 0.7, 6.55, 12, 0.4,
        'Tong cong 18 slides  |  Phuong phap -> Tuong tac -> Cam xuc -> Ket luan',
        sz=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # =================================================================
    # SLIDE 3: QUY TRINH THUC HIEN (PIPELINE OVERVIEW)
    # =================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C_WHITE)
    txt(s, 0.5, 0.3, 12, 0.7, 'QUY TRINH THUC HIEN (PIPELINE)', sz=28, bold=True, color=C_PRIMARY)
    txt(s, 0.5, 0.9, 10, 0.4, 'Toan bo quy trinh tu dong hoa bang Python, gom 6 buoc chinh:', sz=14, color=C_DARK)

    # Pipeline steps - row 1
    step_box(s, 0.3, 1.5, 2.5, 1.3, '1', 'THU THAP', 'Apify TikTok Scraper\n100 videos + comments\ntir @travinhuniversity', C_SECONDARY)
    arrow_right(s, 2.85, 1.95)
    step_box(s, 3.4, 1.5, 2.5, 1.3, '2', 'LAM SACH', 'DataCleaner\nChuan hoa, loai trung\nXoa comment spam', C_TEAL)
    arrow_right(s, 5.95, 1.95)
    step_box(s, 6.5, 1.5, 2.5, 1.3, '3', 'KIEM TRA', 'DataValidator\nKiem tra cau truc\nDanh dau bat thuong', C_WARNING)
    arrow_right(s, 9.05, 1.95)
    step_box(s, 9.6, 1.5, 3.3, 1.3, '4', 'GAN NHAN', 'PhoBERT Sentiment\n3 nhan: Pos/Neu/Neg\nConfidence score', C_PURPLE)

    # Pipeline steps - row 2
    step_box(s, 3.4, 3.3, 2.5, 1.3, '5', 'GHE P NOI', 'Merge video stats\n+ comment sentiment\n=> tong_hop_data', C_ACCENT)
    arrow_right(s, 5.95, 3.75)
    step_box(s, 6.5, 3.3, 2.5, 1.3, '6', 'PHAN TICH', 'Bieu do, thong ke\nEngagement + Sentiment\nXuat bao cao', C_PRIMARY)

    # Technology stack
    txt(s, 0.5, 5.0, 12, 0.5, 'CONG NGHE SU DUNG', sz=16, bold=True, color=C_PRIMARY)

    tech_headers = ['Thanh phan', 'Cong nghe', 'Mo ta']
    tech_rows = [
        ['Thu thap du lieu', 'Apify TikTok Scraper', 'API scraping, khong can dang nhap TikTok'],
        ['Xu ly du lieu', 'Python + Pandas', 'Lam sach, chuan hoa, loai bo trung lap'],
        ['Phan tich cam xuc', 'PhoBERT (wonrax)', 'Model NLP tieng Viet, do chinh xac ~92%'],
        ['Truc quan hoa', 'Matplotlib + Seaborn + WordCloud', '10 bieu do: pie, bar, line, scatter, wordcloud'],
        ['Bao cao', 'python-pptx', 'Tu dong tao PowerPoint 18 slides'],
    ]
    table(s, 0.3, 5.5, 12.7, 1.7, tech_headers, tech_rows, [2.5, 3.5, 6.7])

    # =================================================================
    # SLIDE 4: THU THAP DU LIEU (DATA COLLECTION)
    # =================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C_WHITE)
    txt(s, 0.5, 0.3, 12, 0.7, 'BUOC 1: THU THAP DU LIEU', sz=28, bold=True, color=C_PRIMARY)

    # Left: Apify info
    txt(s, 0.5, 1.1, 6, 0.5, 'Cong cu: Apify TikTok Scraper', sz=18, bold=True, color=C_SECONDARY)
    bullets(s, 0.5, 1.7, 6, 3.5, [
        '> Doi tuong: Kenh @travinhuniversity (DH Tra Vinh)',
        '> So luong thu thap: 100 videos moi nhat',
        '> Ngay thu thap: 06/03/2026',
        '',
        '> Du lieu video bao gom:',
        '  - ID video, ngay dang (create_time)',
        '  - Luot xem (play_count)',
        '  - Luot thich (like_count)',
        '  - Luot binh luan (comment_count)',
        '  - Luot chia se (share_count)',
        '  - Luot luu (collect_count)',
        '',
        '> Du lieu binh luan:',
        '  - 1,222 binh luan tu 100 videos',
        '  - Noi dung, tac gia, so likes',
        '  - Ngon ngu (vi, en, ...)',
    ], sz=12, color=C_DARK)

    # Right: Data structure
    txt(s, 7, 1.1, 5.5, 0.5, 'Cau truc du lieu thu thap', sz=18, bold=True, color=C_SECONDARY)

    struct_headers = ['Truong', 'Kieu', 'Vi du']
    struct_rows = [
        ['id', 'string', '7592439894094482696'],
        ['createTime', 'timestamp', '1717200000'],
        ['stats.playCount', 'int', '2,000,000'],
        ['stats.diggCount', 'int', '122,900'],
        ['stats.commentCount', 'int', '862'],
        ['stats.shareCount', 'int', '26,100'],
        ['stats.collectCount', 'int', '5,200'],
        ['comments[].text', 'string', '"Truong dep qua!"'],
        ['comments[].likes', 'int', '1,454'],
    ]
    table(s, 7, 1.7, 5.8, 3.5, struct_headers, struct_rows, [2, 1, 2.8])

    # Bottom: Format note
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xEB, 0xF5, 0xFB)
    shape.line.color.rgb = C_SECONDARY

    txt(s, 0.8, 5.9, 11.7, 0.3, 'Luu y ve du lieu:', sz=13, bold=True, color=C_SECONDARY)
    txt(s, 0.8, 6.3, 11.7, 0.5,
        'Du lieu duoc xuat tu Apify duoi dang JSON (dataset_tiktok-scraper_2026-03-06.json). '
        'File gom 100 objects, moi object chua thong tin 1 video va danh sach binh luan. '
        'Tong dung luong: ~5MB.',
        sz=11, color=C_DARK)

    # =================================================================
    # SLIDE 5: XU LY & GAN NHAN (PROCESSING & LABELING)
    # =================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C_WHITE)
    txt(s, 0.5, 0.3, 12, 0.7, 'BUOC 2-4: XU LY, KIEM TRA & GAN NHAN', sz=28, bold=True, color=C_PRIMARY)

    # Column 1: Data Cleaning
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.3), Inches(1.1), Inches(4), Inches(5.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xF5)
    shape.line.color.rgb = C_TEAL

    txt(s, 0.5, 1.2, 3.6, 0.5, 'B2: LAM SACH (DataCleaner)', sz=15, bold=True, color=C_TEAL)
    bullets(s, 0.5, 1.8, 3.6, 4.8, [
        '> Chuan hoa Unicode (NFC)',
        '  - Sua loi font chu tieng Viet',
        '  - Thong nhat ky tu dac biet',
        '',
        '> Loai bo binh luan trung lap',
        '  - Kiem tra exact duplicate',
        '  - Xoa comment spam (lien tuc)',
        '',
        '> Lam sach noi dung',
        '  - Xoa comment rong/chi co emoji',
        '  - Giu lai comment co y nghia',
        '',
        '> Chuan hoa thoi gian',
        '  - Unix timestamp -> datetime',
        '  - Luu tru ISO format',
        '',
        '> Ket qua:',
        '  Truoc: 1,500+ binh luan',
        '  Sau: 1,222 binh luan sach',
    ], sz=11, color=C_DARK)

    # Column 2: Data Validation
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(4.6), Inches(1.1), Inches(4), Inches(5.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFE, 0xF9, 0xE7)
    shape.line.color.rgb = C_WARNING

    txt(s, 4.8, 1.2, 3.6, 0.5, 'B3: KIEM TRA (DataValidator)', sz=15, bold=True, color=C_WARNING)
    bullets(s, 4.8, 1.8, 3.6, 4.8, [
        '> Kiem tra cau truc du lieu',
        '  - Verify truong bat buoc ton tai',
        '  - Kiem tra kieu du lieu dung',
        '',
        '> Kiem tra gia tri hop le',
        '  - play_count >= 0',
        '  - like_count >= 0',
        '  - Ngay dang hop le',
        '',
        '> Phat hien bat thuong (outliers)',
        '  - Video co so lieu bat thuong',
        '  - Binh luan co dau hieu bot',
        '',
        '> Danh dau du lieu van de',
        '  - Flag de review thu cong',
        '  - Ghi log bat thuong',
        '',
        '> Ket qua:',
        '  100/100 videos pass validation',
        '  0 bat thuong nghiem trong',
    ], sz=11, color=C_DARK)

    # Column 3: Sentiment Labeling
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(8.9), Inches(1.1), Inches(4.1), Inches(5.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF4, 0xEC, 0xF7)
    shape.line.color.rgb = C_PURPLE

    txt(s, 9.1, 1.2, 3.7, 0.5, 'B4: GAN NHAN (PhoBERT)', sz=15, bold=True, color=C_PURPLE)
    bullets(s, 9.1, 1.8, 3.7, 4.8, [
        '> Model: PhoBERT Vietnamese',
        '  wonrax/phobert-base-',
        '  vietnamese-sentiment',
        '',
        '> 3 nhan cam xuc:',
        '  - Positive (tich cuc)',
        '  - Neutral (trung lap)',
        '  - Negative (tieu cuc)',
        '',
        '> Confidence score (0-1)',
        '  - High: >= 0.9 (50.2%)',
        '  - Medium: 0.7-0.9 (38.0%)',
        '  - Low: < 0.7 (11.9%)',
        '',
        '> Do chinh xac: ~92%',
        '  (tren tap test tieng Viet)',
        '',
        '> Ket qua:',
        '  1,222 binh luan da gan nhan',
        '  523 Pos / 467 Neu / 232 Neg',
    ], sz=11, color=C_DARK)

    # =================================================================
    # SLIDE 6: TONG QUAN - KPI + DONUT
    # =================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C_WHITE)
    txt(s, 0.5, 0.3, 12, 0.7, 'TONG QUAN TUONG TAC', sz=28, bold=True, color=C_PRIMARY)

    kpis = [
        ('10,843,095', 'Total Plays', C_SECONDARY),
        ('367,185', 'Total Likes', C_ACCENT),
        ('1,222', 'Comments (thu thap)', C_SUCCESS),
        ('77,377', 'Total Shares', C_WARNING),
        ('13,353', 'Total Saves', C_PRIMARY),
    ]
    for i, (val, lbl, clr) in enumerate(kpis):
        card(s, 0.3 + i*2.5, 1.1, 2.2, 1.4, val, lbl, clr)

    # Engagement rates
    txt(s, 0.5, 2.8, 5, 0.4, 'Ty le tuong tac trung binh:', sz=16, bold=True, color=C_DARK)
    rates_text = 'Engagement: 3.70%  |  Like: 3.39%  |  Comment: 0.06%  |  Share: 0.71%'
    txt(s, 0.5, 3.3, 8, 0.4, rates_text, sz=14, color=C_DARK)

    # Donut chart
    img(s, 'engagement_donut.png', 7, 2.7, 5.8, 4.5)

    txt(s, 0.5, 4.0, 6.5, 2.5,
        'Nguon: Apify TikTok Scraper | 100 videos\n'
        'Ngay thu thap: 06/03/2026\n\n'
        'Likes chiem 79.1% tong tuong tac.\n'
        'Shares (16.7%) va Saves (2.9%) cho thay\n'
        'noi dung duoc chia se va luu lai.\n\n'
        'Luu y: TikTok bao 6,035 comments tren\n'
        '100 videos, nhung Apify chi thu thap\n'
        'duoc 1,222 binh luan co noi dung.',
        sz=13, color=C_DARK)

    # =================================================================
    # SLIDE 7: BAR CHART - SO SANH TUONG TAC
    # =================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C_WHITE)
    txt(s, 0.5, 0.3, 12, 0.7, 'SO SANH TUONG TAC - TOP 10 VIDEOS', sz=28, bold=True, color=C_PRIMARY)
    img(s, 'video_comparison_bar.png', 0.3, 1.0, 12.7, 6.2)

    # =================================================================
    # SLIDE 8: TOP 10 TABLE + INSIGHTS
    # =================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C_WHITE)
    txt(s, 0.5, 0.3, 12, 0.7, 'TOP 10 VIDEOS THEO LUOT XEM', sz=28, bold=True, color=C_PRIMARY)

    headers = ['#', 'Video ID', 'Plays', 'Likes', 'Comments', 'Shares', 'Eng. Rate']
    rows = [
        ['1', '7592439894094482696', '2,000,000', '122,900', '862', '26,100', '7.72%'],
        ['2', '7594387595551853832', '1,600,000', '32,800', '386', '4,969', '2.42%'],
        ['3', '7546456531819138320', '1,000,000', '42,300', '616', '8,614', '5.27%'],
        ['4', '7553953881168858384', '842,800', '24,100', '354', '8,746', '4.02%'],
        ['5', '7552175330895596807', '482,500', '16,100', '80', '530', '3.57%'],
        ['6', '7551055901289975058', '350,700', '5,135', '189', '935', '1.81%'],
        ['7', '7552179278645841159', '338,900', '6,950', '80', '350', '2.22%'],
        ['8', '7536152035369045249', '233,700', '7,266', '297', '639', '4.10%'],
        ['9', '7573994030976011528', '211,300', '8,971', '30', '922', '4.87%'],
        ['10', '7552350837150862600', '193,400', '4,865', '55', '179', '2.73%'],
    ]
    table(s, 0.3, 1.1, 9, 4.5, headers, rows, [0.4, 2.3, 1.3, 1.1, 1.0, 1.0, 1.0])

    # Insights box
    txt(s, 9.8, 1.1, 3.2, 0.5, 'Nhan xet', sz=18, bold=True, color=C_SECONDARY)
    bullets(s, 9.8, 1.7, 3.3, 4.5, [
        '> Video #1 dat 2M views voi engagement 7.72%',
        '',
        '> Top 4 videos chiem 50% tong views',
        '',
        '> Engagement TB: 3.70%',
        '  (cao hon TB nganh 2-3%)',
        '',
        '> Video nhieu views nhat cung co engagement cao nhat',
        '',
        '> Like chiem phan lon tuong tac, Comment rat thap (0.06%)',
    ], sz=12, color=C_DARK)

    txt(s, 0.3, 6.0, 12, 0.5,
        'Video #1 vuot xa cac video khac: 2M plays, 122.9K likes, 7.72% engagement',
        sz=12, bold=True, color=C_ACCENT)

    # =================================================================
    # SLIDE 9: LINE CHART - XU HUONG
    # =================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C_WHITE)
    txt(s, 0.5, 0.3, 12, 0.7, 'XU HUONG LUOT XEM THEO THOI GIAN', sz=28, bold=True, color=C_PRIMARY)
    img(s, 'views_trend_line.png', 0.3, 1.0, 12.7, 6.2)

    # =================================================================
    # SLIDE 10: SCATTER - PLAYS vs ENGAGEMENT
    # =================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C_WHITE)
    txt(s, 0.5, 0.3, 12, 0.7, 'TUONG QUAN LUOT XEM vs ENGAGEMENT RATE', sz=28, bold=True, color=C_PRIMARY)
    img(s, 'plays_vs_engagement.png', 0.5, 1.0, 12, 6.2)

    # =================================================================
    # SLIDE 11: PIE CHART - SENTIMENT
    # =================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C_WHITE)
    txt(s, 0.5, 0.3, 12, 0.7, 'PHAN TICH SAC THAI BINH LUAN (SENTIMENT)', sz=28, bold=True, color=C_PRIMARY)

    # Pie chart left
    img(s, 'sentiment_pie.png', 0.2, 0.8, 5.5, 5.5)

    # Stats right
    txt(s, 6.2, 1.0, 6.5, 0.5, 'Ket qua Sentiment Analysis', sz=18, bold=True, color=C_SECONDARY)
    txt(s, 6.2, 1.6, 6.5, 0.4, 'Model: PhoBERT (Vietnamese NLP) - Do chinh xac ~92%', sz=13, color=C_DARK)

    # Sentiment details
    conf_headers = ['Sentiment', 'So luong', 'Ty le', 'Avg Conf.']
    conf_rows = [
        ['Positive', '523', '42.8%', '0.815'],
        ['Neutral', '467', '38.2%', '0.931'],
        ['Negative', '232', '19.0%', '0.892'],
    ]
    table(s, 6.2, 2.2, 6.5, 1.5, conf_headers, conf_rows, [1.5, 1.2, 1.2, 1.2])

    # Confidence
    txt(s, 6.2, 4.0, 6.5, 0.4, 'Phan bo Confidence', sz=16, bold=True, color=C_DARK)
    c_headers = ['Muc', 'So luong', 'Ty le']
    c_rows = [
        ['High (>=0.9)', '613', '50.2%'],
        ['Medium (0.7-0.9)', '464', '38.0%'],
        ['Low (<0.7)', '145', '11.9%'],
    ]
    table(s, 6.2, 4.5, 6.5, 1.2, c_headers, c_rows, [2.5, 2, 2])

    txt(s, 6.2, 6.0, 6.5, 1,
        'Nhan xet: 42.8% binh luan tich cuc - Kha quan cho hinh anh TVU. '
        '50.2% co confidence >= 0.9 cho thay model phan tich dang tin cay.',
        sz=12, bold=True, color=C_DARK)

    # =================================================================
    # SLIDE 12: SENTIMENT CHARTS (confidence + video sentiment)
    # =================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C_WHITE)
    txt(s, 0.5, 0.3, 12, 0.7, 'CHI TIET SENTIMENT - CONFIDENCE & THEO VIDEO', sz=28, bold=True, color=C_PRIMARY)
    img(s, 'confidence_hist.png', 0.2, 1.0, 6.3, 3.5)
    img(s, 'video_sentiment_stack.png', 6.8, 1.0, 6.3, 3.5)

    # Highlights
    txt(s, 0.5, 4.8, 6, 0.4, 'Nhan xet Confidence:', sz=14, bold=True, color=C_SECONDARY)
    bullets(s, 0.5, 5.3, 6, 1.5, [
        '> Neutral co confidence cao nhat (0.931)',
        '> 88.1% binh luan co confidence >= 0.7',
        '> Chi 11.9% co confidence thap (<0.7)',
    ], sz=12, color=C_DARK)

    txt(s, 7, 4.8, 6, 0.4, 'Nhan xet theo Video:', sz=14, bold=True, color=C_SECONDARY)
    bullets(s, 7, 5.3, 6, 1.5, [
        '> Da so video co binh luan Positive nhieu hon',
        '> Mot so video co nhieu binh luan Negative',
        '  => Can xem xet noi dung de cai thien',
    ], sz=12, color=C_DARK)

    # =================================================================
    # SLIDE 13: WORD CLOUD
    # =================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C_WHITE)
    txt(s, 0.5, 0.3, 12, 0.7, 'TU KHOA NOI BAT TRONG BINH LUAN', sz=28, bold=True, color=C_PRIMARY)
    img(s, 'wordcloud.png', 0.3, 1.0, 12.7, 6.2)

    # =================================================================
    # SLIDE 14: WORD CLOUD TICH CUC vs TIEU CUC
    # =================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C_WHITE)
    txt(s, 0.5, 0.3, 12, 0.7, 'TU KHOA - TICH CUC vs TIEU CUC', sz=28, bold=True, color=C_PRIMARY)

    txt(s, 0.5, 1.0, 6, 0.5, 'Binh luan TICH CUC (523)', sz=16, bold=True, color=C_SUCCESS)
    img(s, 'wordcloud_positive.png', 0.2, 1.5, 6.3, 3.8)

    txt(s, 7, 1.0, 6, 0.5, 'Binh luan TIEU CUC (232)', sz=16, bold=True, color=C_ACCENT)
    img(s, 'wordcloud_negative.png', 6.8, 1.5, 6.3, 3.8)

    bullets(s, 0.5, 5.5, 5.5, 1.5, [
        '> Tu khoa tich cuc: tu hao, yeu thuong, dep, gioi...',
        '> Noi dung truyen cam hung, khien nguoi xem cam thay gan gui',
    ], sz=12, color=C_SUCCESS)

    bullets(s, 7, 5.5, 5.5, 1.5, [
        '> Tu khoa tieu cuc: gop y, phan anh, thieu thong tin...',
        '> Phan lon la gop y xay dung, khong co noi dung doc hai',
    ], sz=12, color=C_ACCENT)

    # =================================================================
    # SLIDE 15: TOP COMMENTS
    # =================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C_WHITE)
    txt(s, 0.5, 0.3, 12, 0.7, 'TOP BINH LUAN TIEU BIEU', sz=28, bold=True, color=C_PRIMARY)

    txt(s, 0.3, 1.0, 6, 0.4, 'Top 5 Positive Comments (nhieu likes)', sz=15, bold=True, color=C_SUCCESS)
    pos_h = ['#', 'Likes', 'Conf', 'Noi dung']
    pos_r = [
        ['1', '1,454', '0.85', 'O Ha Lan duoc 6 nam lang le doc cmt va cuoi'],
        ['2', '982', '0.82', 'nay moi la chu ne con t la dia rau muong xao'],
        ['3', '774', '0.73', 'Da qua z. Thich truong nao co CT trao doi SV...'],
        ['4', '386', '0.87', 'DH Can Tho ma que Tra Vinh chao nhau cai duoc khong'],
        ['5', '251', '0.74', '3 nam nhu 1'],
    ]
    table(s, 0.3, 1.5, 6.2, 2.2, pos_h, pos_r, [0.3, 0.7, 0.6, 4.6])

    txt(s, 0.3, 4.0, 6, 0.4, 'Top 5 Negative Comments (nhieu likes)', sz=15, bold=True, color=C_ACCENT)
    neg_h = ['#', 'Likes', 'Conf', 'Noi dung']
    neg_r = [
        ['1', '2,093', '0.98', 'Theo gv tieu hoc ma chu nhu ve bua'],
        ['2', '190', '0.97', 'nghe Dieu Kien hat xong thi thay Duong ai Vy...'],
        ['3', '187', '0.58', 'Vay la thang be nay nho hon minh 10t saoo'],
        ['4', '142', '0.97', 'Khong nhu may chi em nghi dau, no hon loan...'],
        ['5', '92', '0.59', 'sao may dai hoc tren SG ko co hic'],
    ]
    table(s, 0.3, 4.5, 6.2, 2.2, neg_h, neg_r, [0.3, 0.7, 0.6, 4.6])

    # Right side insights
    txt(s, 7, 1.0, 5.5, 0.5, 'Phan tich & Nhan xet', sz=18, bold=True, color=C_SECONDARY)
    bullets(s, 7, 1.6, 5.5, 5.5, [
        '> Binh luan Positive:',
        '  - Tu hao ve truong, yeu thuong truong',
        '  - Khen ngoi hoat dong, su kien',
        '  - Hoc vien/cuu sv gop cam xuc tich cuc',
        '',
        '> Binh luan Negative:',
        '  - Gop y chinh dang (chat luong giang day)',
        '  - Phan anh chinh sach, quy dinh',
        '  - Khong co noi dung doc hai',
        '',
        '> Dac biet:',
        '  - Comment tieu cuc #1 co 2,093 likes',
        '    (nhieu nhat trong tat ca comments)',
        '  - Can chu y xu ly noi dung nay',
        '',
        '> De xuat:',
        '  - Tang noi dung truyen cam hung',
        '  - Tra loi binh luan negative kip thoi',
        '  - Tao noi dung tuong tac voi sinh vien',
    ], sz=12, color=C_DARK)

    # =================================================================
    # SLIDE 16: KET HOP ENGAGEMENT x SENTIMENT
    # =================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C_WHITE)
    txt(s, 0.5, 0.3, 12, 0.7, 'PHAN TICH KET HOP: ENGAGEMENT x SENTIMENT', sz=28, bold=True, color=C_PRIMARY)

    txt(s, 0.3, 1.0, 6, 0.4, 'Videos TICH CUC (High Engagement + Positive)', sz=14, bold=True, color=C_SUCCESS)
    good_h = ['#', 'Video ID', 'Plays', 'Eng.', 'Pos', 'Neg', 'Score']
    good_r = [
        ['1', '7581690429377187079', '41,100', '1.28%', '2', '0', '1.00'],
        ['2', '7581018792151502100', '79,100', '1.90%', '2', '0', '1.00'],
        ['3', '7573623876056419604', '28,600', '2.50%', '1', '0', '1.00'],
        ['4', '7573356202231352583', '20,900', '2.94%', '1', '0', '1.00'],
        ['5', '7543952673611992328', '9,853', '1.13%', '1', '0', '1.00'],
    ]
    table(s, 0.3, 1.5, 6.2, 2, good_h, good_r, [0.3, 2.2, 0.8, 0.7, 0.5, 0.5, 0.7])

    txt(s, 0.3, 3.8, 6, 0.4, 'Videos CAN QUAN TAM (High Engagement + Negative)', sz=14, bold=True, color=C_ACCENT)
    bad_h = ['#', 'Video ID', 'Plays', 'Eng.', 'Pos', 'Neg', 'Score']
    bad_r = [
        ['1', '7606175475567824136', '13,100', '5.37%', '0', '4', '-1.00'],
        ['2', '7573617332363201799', '8,041', '1.27%', '0', '1', '-1.00'],
        ['3', '7599581460130958609', '44,500', '3.92%', '0', '1', '-0.50'],
        ['4', '7586154838279326993', '16,100', '2.78%', '0', '1', '-0.50'],
        ['5', '7605242871989685522', '31,400', '3.09%', '0', '2', '-0.33'],
    ]
    table(s, 0.3, 4.3, 6.2, 2, bad_h, bad_r, [0.3, 2.2, 0.8, 0.7, 0.5, 0.5, 0.7])

    # Right insights
    txt(s, 7, 1.0, 5.5, 0.5, 'Y nghia', sz=18, bold=True, color=C_SECONDARY)
    bullets(s, 7, 1.6, 5.5, 5, [
        '> Score = (Pos - Neg) / (Pos + Neg + Neutral)',
        '  Score +1.0 = 100% tich cuc',
        '  Score -1.0 = 100% tieu cuc',
        '',
        '> Videos TICH CUC (score > 0.5):',
        '  - Noi dung hap dan, phan hoi tot',
        '  - Nen lam them noi dung tuong tu',
        '',
        '> Videos CAN QUAN TAM (score < 0):',
        '  - Co nhieu binh luan tieu cuc',
        '  - Can review noi dung video',
        '  - Tra loi/giai dap de cai thien',
        '',
        '> Video 7606175... co engagement 5.37%',
        '  nhung 4 binh luan deu tieu cuc',
        '  => Engagement cao khong phai luc nao',
        '     cung la dau hieu tot',
    ], sz=12, color=C_DARK)

    # =================================================================
    # SLIDE 17: TOM TAT & DE XUAT
    # =================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C_WHITE)
    txt(s, 0.5, 0.3, 12, 0.7, 'TOM TAT & DE XUAT', sz=28, bold=True, color=C_PRIMARY)

    # Key findings
    txt(s, 0.5, 1.1, 6, 0.5, 'Ket qua chinh', sz=20, bold=True, color=C_SECONDARY)
    bullets(s, 0.5, 1.7, 6, 4, [
        '> 100 videos, 10.8M luot xem, 1,222 binh luan',
        '> Engagement Rate TB: 3.70% (cao hon TB nganh)',
        '> Top video: 2M views, 122.9K likes',
        '',
        '> Sentiment: 42.8% Positive, 19.0% Negative',
        '> PhoBERT accuracy: ~92% cho tieng Viet',
        '> 88.1% binh luan co confidence >= 0.7',
        '',
        '> Likes chiem 79.1% tong tuong tac',
        '> Comment rate thap (0.06%)',
        '> Share rate kha (0.71%)',
    ], sz=13, color=C_DARK)

    # Recommendations
    txt(s, 7, 1.1, 5.5, 0.5, 'De xuat cai thien', sz=20, bold=True, color=C_ACCENT)
    bullets(s, 7, 1.7, 5.5, 4, [
        '> Tang tuong tac Comment:',
        '  - Dat cau hoi trong video/caption',
        '  - Tao noi dung debate/chia se y kien',
        '',
        '> Xu ly binh luan tieu cuc:',
        '  - Phan hoi kip thoi binh luan negative',
        '  - Review video co nhieu BL tieu cuc',
        '',
        '> Noi dung nen lam them:',
        '  - Video truyen cam hung (doi song SV)',
        '  - Video tuong tac (Q&A, talkshow)',
        '  - Noi dung viral (trend, thach thuc)',
    ], sz=13, color=C_DARK)

    # Bottom highlight
    shape = s.shapes.add_shape(1, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_PRIMARY
    shape.line.fill.background()
    txt(s, 0.7, 6.3, 12, 0.5,
        'KET LUAN: Kenh TikTok TVU co hieu suat tot (3.70% engagement). '
        'Can tang tuong tac comment va xu ly binh luan tieu cuc de cai thien hinh anh.',
        sz=14, bold=True, color=C_WHITE)

    # =================================================================
    # SLIDE 18: THANK YOU
    # =================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C_PRIMARY)
    txt(s, 1, 1.5, 11, 1.2, 'CAM ON DA LANG NGHE!', sz=44, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(s, 1, 3.0, 11, 0.8, 'Q&A - Thac mac va gop y', sz=24, color=C_LIGHT, align=PP_ALIGN.CENTER)

    # Divider
    shape = s.shapes.add_shape(1, Inches(4), Inches(4.0), Inches(5.33), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_LIGHT
    shape.line.fill.background()

    txt(s, 1, 4.3, 11, 0.5, '@travinhuniversity | TikTok Analytics', sz=16, color=C_LIGHT, align=PP_ALIGN.CENTER)
    txt(s, 1, 5.0, 11, 0.5, 'Cong cu: Apify + PhoBERT + Python Pipeline', sz=13, color=C_GRAY, align=PP_ALIGN.CENTER)
    txt(s, 1, 5.5, 11, 0.5, 'Truong Dai hoc Tra Vinh - Bao cao tuan 03/03 - 10/03/2026', sz=13, color=C_GRAY, align=PP_ALIGN.CENTER)

    # Save
    out = Path('reports') / 'TikTok_Analytics_TVU_2026_v5.pptx'
    prs.save(str(out))
    print(f'[OK] PowerPoint saved: {out}  ({len(prs.slides)} slides)')
    return str(out)


if __name__ == '__main__':
    create()
